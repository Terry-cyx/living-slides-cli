"""PPTX → HTML adopter (Phase D1).

Lets users bring an existing PowerPoint file into the iteration loop.
The conversion is intentionally lossy: shapes, animations, master
slides, and exotic geometry are dropped. What survives is text content,
images, and slide order — which is what the cascade needs to be useful.

Adapted from `frontend-slides/scripts/extract-pptx.py` (MIT) by
zarazhangrui. Image extraction logic, shape iteration, and notes
handling all originate there.
"""

from __future__ import annotations

import html as html_lib
from pathlib import Path


_MISSING_DEP_MESSAGE = (
    "PPTX adoption requires python-pptx. Install it with:\n"
    "    uv tool install 'living-slides[pptx]'\n"
    "or:\n"
    "    pip install 'living-slides[pptx]'"
)


def _require_pptx():
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError as e:
        raise RuntimeError(_MISSING_DEP_MESSAGE) from e


def adopt_pptx(pptx_path: str, output_html: str | None = None) -> str:
    """Convert a PPTX file to a slive HTML deck and write it next to the source.

    Returns the path to the generated HTML file. Images are extracted to
    `<deck>-assets/` next to the output HTML, matching the convention used
    by `living_slides.assets`.
    """
    _require_pptx()
    from pptx import Presentation

    src = Path(pptx_path)
    if not src.exists():
        raise FileNotFoundError(f"PPTX file not found: {src}")

    if output_html is None:
        output_html = str(src.with_suffix(".html"))
    out_path = Path(output_html)
    assets_dir = out_path.parent / f"{out_path.stem}-assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(src))
    slides_html: list[str] = []

    for slide_index, slide in enumerate(prs.slides):
        slides_html.append(_render_slide(slide, slide_index, assets_dir))

    title = src.stem.replace("-", " ").replace("_", " ").title()
    document = _wrap_document(title, "\n".join(slides_html))
    out_path.write_text(document, encoding="utf-8")
    return str(out_path)


def _render_slide(slide, slide_index: int, assets_dir: Path) -> str:
    title_text = ""
    body_chunks: list[str] = []
    image_chunks: list[str] = []
    notes_text = ""

    title_shape = None
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape is title_shape and not title_text:
                title_text = text
            else:
                body_chunks.append(text)

        # shape_type 13 = PICTURE — preserved from frontend-slides original.
        if getattr(shape, "shape_type", None) == 13:
            try:
                image = shape.image
                ext = image.ext
                fname = f"slide{slide_index + 1}_img{len(image_chunks) + 1}.{ext}"
                (assets_dir / fname).write_bytes(image.blob)
                rel = f"{assets_dir.name}/{fname}"
                image_chunks.append(
                    f'    <img data-oid="s{slide_index}-image-{len(image_chunks) + 1}" '
                    f'src="{html_lib.escape(rel, quote=True)}" alt="">'
                )
            except Exception:
                # Some shape types report as PICTURE but lack a usable .image
                continue

    if slide.has_notes_slide:
        try:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes_text = ""

    parts: list[str] = []
    parts.append(
        f'<section class="slide" data-oid="s{slide_index}-slide" id="slide-{slide_index + 1}">'
    )
    if title_text:
        parts.append(
            f'    <h2 data-oid="s{slide_index}-title">{html_lib.escape(title_text)}</h2>'
        )
    for i, chunk in enumerate(body_chunks):
        # PPTX text frames frequently contain newlines per paragraph; render
        # each line as its own <p> so the cascade can act on individual
        # statements rather than a single mega-blob.
        for j, line in enumerate(chunk.splitlines()):
            line = line.strip()
            if not line:
                continue
            parts.append(
                f'    <p data-oid="s{slide_index}-text-{i + 1}-{j + 1}">'
                f'{html_lib.escape(line)}</p>'
            )
    parts.extend(image_chunks)
    if notes_text:
        parts.append(
            f'    <aside data-oid="s{slide_index}-notes" hidden>'
            f'{html_lib.escape(notes_text)}</aside>'
        )
    parts.append("</section>")
    return "\n".join(parts)


def _wrap_document(title: str, slides_html: str) -> str:
    return f"""\
<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{html_lib.escape(title)}</title>
<style>
  body {{
    margin: 0;
    background: #0b0d12;
    color: #f3f4f6;
    font-family: -apple-system, "Segoe UI", Inter, sans-serif;
    line-height: 1.5;
  }}
  section.slide {{
    width: 100vw;
    min-height: 100vh;
    padding: clamp(2rem, 6vw, 6rem);
    box-sizing: border-box;
    display: flex;
    flex-direction: column;
    justify-content: center;
    border-bottom: 1px solid #1f2937;
  }}
  h2 {{ font-size: clamp(1.75rem, 4vw, 3rem); margin: 0 0 1rem; }}
  p {{ font-size: clamp(1rem, 1.5vw, 1.25rem); max-width: 60ch; }}
  img {{ max-width: 100%; height: auto; }}
</style>
</head>
<body>
{slides_html}
</body>
</html>
"""
