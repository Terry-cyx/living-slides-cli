"""Tests for the PPTX adopter (Phase D1)."""

from __future__ import annotations

import re
from pathlib import Path

import pytest

# Skip the entire module if python-pptx isn't installed — it's an optional extra.
pptx = pytest.importorskip("pptx")

from pptx.util import Inches  # noqa: E402

from living_slides.adopters.pptx import adopt_pptx  # noqa: E402


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """Build a 3-slide PPTX on the fly."""
    prs = pptx.Presentation()
    layout_title = prs.slide_layouts[0]
    layout_content = prs.slide_layouts[1]

    # Slide 1: title only
    s1 = prs.slides.add_slide(layout_title)
    s1.shapes.title.text = "Cover Slide"
    if len(s1.placeholders) > 1:
        s1.placeholders[1].text = "Subtitle line"

    # Slide 2: title + bullets
    s2 = prs.slides.add_slide(layout_content)
    s2.shapes.title.text = "Agenda"
    body = s2.placeholders[1].text_frame
    body.text = "First point"
    body.add_paragraph().text = "Second point"
    body.add_paragraph().text = "Third point"

    # Slide 3: title + speaker notes
    s3 = prs.slides.add_slide(layout_title)
    s3.shapes.title.text = "Done"
    s3.notes_slide.notes_text_frame.text = "Remember to thank the audience"

    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    return out


def test_adopt_pptx_writes_html_next_to_source(sample_pptx):
    out = adopt_pptx(str(sample_pptx))
    assert Path(out).exists()
    assert Path(out).suffix == ".html"


def test_adopt_pptx_renders_three_slides(sample_pptx):
    out = adopt_pptx(str(sample_pptx))
    html = Path(out).read_text(encoding="utf-8")
    assert html.count('class="slide"') == 3


def test_adopt_pptx_carries_data_oid_on_every_slide(sample_pptx):
    out = adopt_pptx(str(sample_pptx))
    html = Path(out).read_text(encoding="utf-8")
    # Each slide section gets a unique slide oid.
    section_oids = re.findall(r'<section[^>]*data-oid="([^"]+)"', html)
    assert len(section_oids) == 3
    assert len(set(section_oids)) == 3


def test_adopt_pptx_preserves_titles_and_body(sample_pptx):
    out = adopt_pptx(str(sample_pptx))
    html = Path(out).read_text(encoding="utf-8")
    assert "Cover Slide" in html
    assert "Agenda" in html
    assert "First point" in html
    assert "Second point" in html
    assert "Done" in html


def test_adopt_pptx_includes_speaker_notes(sample_pptx):
    out = adopt_pptx(str(sample_pptx))
    html = Path(out).read_text(encoding="utf-8")
    assert "Remember to thank the audience" in html


def test_adopt_pptx_custom_output_path(sample_pptx, tmp_path):
    target = tmp_path / "elsewhere" / "out.html"
    target.parent.mkdir()
    result = adopt_pptx(str(sample_pptx), output_html=str(target))
    assert Path(result) == target
    assert target.exists()


def test_adopt_pptx_missing_file_raises(tmp_path):
    with pytest.raises(FileNotFoundError):
        adopt_pptx(str(tmp_path / "nope.pptx"))


# ---- CLI integration ----

def test_cli_adopt_from_pptx(sample_pptx, tmp_path):
    from click.testing import CliRunner
    from living_slides.cli import main

    runner = CliRunner()
    result = runner.invoke(main, ["adopt", str(sample_pptx), "--from", "pptx"])
    assert result.exit_code == 0, result.output
    expected = sample_pptx.with_suffix(".html")
    assert expected.exists()
    assert "data-oid" in expected.read_text(encoding="utf-8")
